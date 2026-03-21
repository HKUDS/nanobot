"""PowerPoint (.pptx) reading and analysis tools.

Level 1 — read_pptx: text/notes/tables extraction (python-pptx only)
Level 2 — analyze_pptx: optional slide rendering (LibreOffice + pdftoppm) +
          per-slide LLM vision analysis
Level 3 — analyze_pptx: deck-level synthesis (executive summary, risks,
          decisions, action items, unanswered questions)
Level 4 — pptx_get_slide: cached slide retrieval for follow-up Q&A
"""

from __future__ import annotations

import asyncio
import base64
import json
import shutil
import tempfile
from pathlib import Path
from typing import TYPE_CHECKING, Any

from nanobot.agent.tools.base import Tool, ToolResult
from nanobot.agent.tools.filesystem import _resolve_path

if TYPE_CHECKING:
    from nanobot.agent.tools.result_cache import ToolResultCache

# ---------------------------------------------------------------------------
# Constants & prompts
# ---------------------------------------------------------------------------

_DEFAULT_VISION_MODEL = "gpt-4o-mini"
MAX_CONCURRENT = 5

SLIDE_ANALYSIS_PROMPT = """\
You are analyzing one PowerPoint slide.
You are given the extracted text content (JSON) and optionally a rendered slide image.
Use BOTH the text and image (if provided) for your analysis.

Return a JSON object with these keys:
- title: string
- summary: string (1-3 sentences)
- key_points: string[] (main points on this slide)
- decisions: string[] (decisions mentioned or implied)
- risks: string[] (risks, concerns, blockers)
- action_items: string[] (tasks, follow-ups, to-dos)
- deadlines: string[] (dates, timelines, milestones)
- owners: string[] (people, teams, roles responsible)
- chart_insights: string[] (what charts/graphs show)
- visual_observations: string[] (layout, emphasis, diagrams, screenshots)

Omit keys with empty arrays. Be specific and cite actual content from the slide.\
"""

DECK_SYNTHESIS_PROMPT = """\
You are synthesizing a complete PowerPoint deck analysis.
You are given per-slide analyses as a JSON array.

Return a JSON object with these keys:
- executive_summary: string (concise 2-4 paragraph overview of the entire deck)
- risks: string[] (all risks across the deck, with slide numbers, deduplicated)
- decisions: string[] (all decisions, with slide numbers)
- action_items: string[] (all action items, include owners and deadlines where known)
- deadlines: string[] (all deadlines and timelines mentioned)
- unanswered_questions: string[] (gaps, unclear points, missing information)
- themes: string[] (recurring themes across the deck)

Be thorough. Always cite slide numbers. Deduplicate across slides.\
"""


# ---------------------------------------------------------------------------
# Helpers — text extraction
# ---------------------------------------------------------------------------


def _import_pptx() -> Any:
    """Import python-pptx Presentation class with helpful error."""
    try:
        from pptx import Presentation  # type: ignore[import-untyped]

        return Presentation
    except ImportError as err:
        raise ImportError(
            "python-pptx is required for PowerPoint reading. "
            "Install it with: pip install 'nanobot-ai[pptx]'"
        ) from err


def _extract_slides_data(file_path: Path) -> list[dict[str, Any]]:
    """Extract structured data from each slide as a list of dicts."""
    presentation_cls = _import_pptx()
    prs = presentation_cls(str(file_path))
    slides: list[dict[str, Any]] = []

    for idx, slide in enumerate(prs.slides, start=1):
        title = ""
        text_blocks: list[str] = []
        tables: list[list[list[str]]] = []
        shape_types: list[str] = []

        for shape in slide.shapes:
            shape_types.append(type(shape).__name__)

            text = getattr(shape, "text", "")
            if isinstance(text, str) and text.strip():
                text = text.strip()
                text_blocks.append(text)
                if not title:
                    title = text.splitlines()[0][:180]

            if getattr(shape, "has_table", False):
                rows: list[list[str]] = []
                for row in shape.table.rows:
                    rows.append([cell.text.strip() for cell in row.cells])
                if rows:
                    tables.append(rows)

        notes = ""
        if getattr(slide, "has_notes_slide", False):
            note_chunks: list[str] = []
            for shape in slide.notes_slide.shapes:
                t = getattr(shape, "text", "")
                if isinstance(t, str) and t.strip():
                    note_chunks.append(t.strip())
            notes = "\n".join(note_chunks)

        slides.append(
            {
                "slide_number": idx,
                "title": title,
                "text_blocks": text_blocks,
                "tables": tables,
                "notes": notes,
                "shape_types": shape_types,
            }
        )
    return slides


def _format_slides_markdown(file_name: str, slides: list[dict[str, Any]]) -> str:
    """Format structured slide data as readable markdown."""
    sections: list[str] = []
    sections.append(f"# PowerPoint: {file_name} ({len(slides)} slides)\n")

    for slide in slides:
        idx = slide["slide_number"]
        title = slide.get("title", "")
        section = f"## Slide {idx}"
        if title:
            section += f": {title}"
        section += "\n"

        text_blocks = slide.get("text_blocks", [])
        if text_blocks:
            section += "\n### Content\n"
            section += "\n\n".join(text_blocks) + "\n"

        tables = slide.get("tables", [])
        if tables:
            section += "\n### Tables\n"
            for i, table in enumerate(tables, start=1):
                if len(tables) > 1:
                    section += f"\nTable {i}:\n"
                for row in table:
                    section += " | ".join(row) + "\n"

        notes = slide.get("notes", "")
        if notes:
            section += "\n### Speaker Notes\n"
            section += notes + "\n"

        shape_types = slide.get("shape_types", [])
        non_text = [s for s in shape_types if s not in ("Shape",)]
        if non_text:
            section += f"\n*Shapes: {', '.join(non_text)}*\n"

        sections.append(section)

    return "\n".join(sections)


# ---------------------------------------------------------------------------
# Helpers — slide rendering
# ---------------------------------------------------------------------------


async def _render_slides(pptx_path: Path, out_dir: Path) -> list[Path] | None:
    """Render slides to PNG via LibreOffice + pdftoppm.

    Returns a sorted list of PNG paths, or ``None`` if rendering is unavailable.
    """
    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    pdftoppm = shutil.which("pdftoppm")
    if not soffice or not pdftoppm:
        return None

    out_dir.mkdir(parents=True, exist_ok=True)

    # pptx → PDF
    proc = await asyncio.create_subprocess_exec(
        soffice,
        "--headless",
        "--convert-to",
        "pdf",
        "--outdir",
        str(out_dir),
        str(pptx_path),
        stdout=asyncio.subprocess.PIPE,
        stderr=asyncio.subprocess.PIPE,
    )
    await proc.communicate()
    if proc.returncode != 0:
        return None

    pdf_path = out_dir / (pptx_path.stem + ".pdf")
    if not pdf_path.exists():
        return None

    # PDF → PNGs (150 dpi for good quality without huge files)
    prefix = str(out_dir / "slide")
    proc = await asyncio.create_subprocess_exec(
        pdftoppm,
        "-png",
        "-r",
        "150",
        str(pdf_path),
        prefix,
        stdout=asyncio.subprocess.PIPE,
        stderr=asyncio.subprocess.PIPE,
    )
    await proc.communicate()
    if proc.returncode != 0:
        return None

    images = sorted(out_dir.glob("slide-*.png"))
    return images if images else None


# ---------------------------------------------------------------------------
# Helpers — LLM analysis
# ---------------------------------------------------------------------------


def _encode_image_base64(path: Path) -> str:
    """Base64-encode an image file."""
    return base64.b64encode(path.read_bytes()).decode("ascii")


async def _call_llm(messages: list[dict[str, Any]], model: str) -> str:
    """Call LLM via litellm and return response text."""
    import litellm

    response = await litellm.acompletion(model=model, messages=messages, temperature=0.2)
    return response.choices[0].message.content or ""


def _parse_json_response(text: str) -> dict[str, Any]:
    """Parse JSON from LLM response, stripping markdown code fences."""
    cleaned = text.strip()
    if cleaned.startswith("```"):
        lines = cleaned.split("\n")
        lines = [line for line in lines[1:] if not line.strip().startswith("```")]
        cleaned = "\n".join(lines)
    try:
        result: dict[str, Any] = json.loads(cleaned)
        return result
    except json.JSONDecodeError:
        return {"raw_response": text}


async def _analyze_slide(
    slide_data: dict[str, Any],
    image_path: Path | None,
    model: str,
) -> dict[str, Any]:
    """Analyze one slide via LLM (with vision if image available)."""
    content: list[dict[str, Any]] = [
        {"type": "text", "text": json.dumps(slide_data, ensure_ascii=False)},
    ]
    if image_path and image_path.exists():
        b64 = _encode_image_base64(image_path)
        content.append(
            {
                "type": "image_url",
                "image_url": {"url": f"data:image/png;base64,{b64}"},
            }
        )

    messages: list[dict[str, Any]] = [
        {"role": "system", "content": SLIDE_ANALYSIS_PROMPT},
        {"role": "user", "content": content},
    ]
    text = await _call_llm(messages, model)
    analysis = _parse_json_response(text)

    # Preserve source grounding
    analysis["slide_number"] = slide_data["slide_number"]
    analysis["source_title"] = slide_data.get("title", "")
    analysis["source_text_blocks"] = slide_data.get("text_blocks", [])
    analysis["source_notes"] = slide_data.get("notes", "")
    return analysis


async def _synthesize_deck(
    slide_analyses: list[dict[str, Any]],
    model: str,
) -> dict[str, Any]:
    """Produce deck-level synthesis from per-slide analyses."""
    messages: list[dict[str, Any]] = [
        {"role": "system", "content": DECK_SYNTHESIS_PROMPT},
        {"role": "user", "content": json.dumps(slide_analyses, ensure_ascii=False)},
    ]
    text = await _call_llm(messages, model)
    return _parse_json_response(text)


def _format_analysis_output(
    file_name: str,
    total_slides: int,
    mode: str,
    synthesis: dict[str, Any],
    output_path: Path,
) -> str:
    """Format analysis as readable markdown for the agent."""
    lines: list[str] = [
        f"# Analysis: {file_name} ({total_slides} slides, {mode} mode)\n",
    ]

    es = synthesis.get("executive_summary", "")
    if es:
        lines.append(f"## Executive Summary\n\n{es}\n")

    for key, heading in [
        ("risks", "Risks"),
        ("decisions", "Decisions"),
        ("action_items", "Action Items"),
        ("deadlines", "Deadlines"),
        ("unanswered_questions", "Unanswered Questions"),
        ("themes", "Themes"),
    ]:
        items = synthesis.get(key, [])
        if items:
            lines.append(f"## {heading}\n")
            for item in items:
                lines.append(f"- {item}")
            lines.append("")

    lines.append(f"\n---\nFull analysis saved to: {output_path}")
    lines.append("Use `read_file` on the analysis JSON for detailed per-slide data.")
    return "\n".join(lines)


# ---------------------------------------------------------------------------
# Tool: read_pptx (Level 1)
# ---------------------------------------------------------------------------


class ReadPptxTool(Tool):
    """Extract structured text content from a PowerPoint (.pptx) file.

    Each slide is cached individually so downstream tools
    (``pptx_get_slide``) can retrieve slide data without the entire
    presentation needing to fit in context.  The tool returns a compact
    metadata envelope (slide count, titles, per-slide cache keys) — never
    raw slide content.
    """

    readonly = True
    cacheable = False

    def __init__(
        self,
        workspace: Path | None = None,
        allowed_dir: Path | None = None,
        cache: ToolResultCache | None = None,
    ):
        self._workspace = workspace
        self._allowed_dir = allowed_dir
        self._cache: ToolResultCache | None = cache

    @property
    def name(self) -> str:
        return "read_pptx"

    @property
    def description(self) -> str:
        return (
            "Read a PowerPoint (.pptx) file and extract all slide content: "
            "titles, text, speaker notes, and tables. Returns metadata with "
            "per-slide cache keys. Use pptx_get_slide to retrieve individual "
            "slide content."
        )

    @property
    def parameters(self) -> dict[str, Any]:
        return {
            "type": "object",
            "properties": {
                "path": {
                    "type": "string",
                    "description": "Path to the .pptx file to read.",
                },
            },
            "required": ["path"],
        }

    def _cache_slide(self, file_path: str, slide_number: int, slide_data: dict[str, Any]) -> str:
        """Cache a single slide's data and return the cache key."""
        if not self._cache:
            return ""
        full_json = json.dumps(slide_data, ensure_ascii=False, default=str)
        return self._cache.store(
            "read_pptx",
            {"path": file_path, "slide": slide_number},
            full_json,
            "",
            token_estimate=len(full_json) // 4,
        )

    async def execute(self, path: str, **kwargs: Any) -> ToolResult:  # type: ignore[override]
        try:
            file_path = _resolve_path(str(path), self._workspace, self._allowed_dir)
        except PermissionError as exc:
            return ToolResult.fail(str(exc), error_type="permission_denied")

        if not file_path.exists():
            return ToolResult.fail(f"File not found: {path}", error_type="not_found")
        if not file_path.is_file():
            return ToolResult.fail(f"Not a file: {path}", error_type="invalid_path")
        if file_path.suffix.lower() != ".pptx":
            return ToolResult.fail(f"Not a .pptx file: {path}", error_type="invalid_format")

        try:
            slides = _extract_slides_data(file_path)
        except ImportError as e:
            return ToolResult.fail(str(e), error_type="missing_dependency")
        except Exception as e:
            return ToolResult.fail(
                f"Error reading PowerPoint file: {e}", error_type="extraction_error"
            )

        if not slides:
            return ToolResult.ok("The presentation has no slides.")

        # Cache each slide and build metadata envelope
        slide_meta: list[dict[str, Any]] = []
        for slide in slides:
            cache_key = self._cache_slide(str(file_path), slide["slide_number"], slide)
            entry: dict[str, Any] = {
                "slide_number": slide["slide_number"],
                "title": slide.get("title", ""),
                "has_tables": bool(slide.get("tables")),
                "has_notes": bool(slide.get("notes")),
            }
            if cache_key:
                entry["cache_key"] = cache_key
            slide_meta.append(entry)

        result = {
            "file": file_path.name,
            "total_slides": len(slides),
            "slides": slide_meta,
        }
        return ToolResult.ok(json.dumps(result, ensure_ascii=False))


# ---------------------------------------------------------------------------
# Tool: analyze_pptx (Level 2+3)
# ---------------------------------------------------------------------------


class AnalyzePptxTool(Tool):
    """Deep multimodal analysis of a PowerPoint deck using LLM."""

    readonly = False  # writes analysis JSON + rendered images

    def __init__(
        self,
        workspace: Path | None = None,
        allowed_dir: Path | None = None,
        cache: ToolResultCache | None = None,
        vision_model: str = _DEFAULT_VISION_MODEL,
    ):
        self._workspace = workspace
        self._allowed_dir = allowed_dir
        self._cache: ToolResultCache | None = cache
        self._vision_model = vision_model

    @property
    def name(self) -> str:
        return "analyze_pptx"

    @property
    def description(self) -> str:
        return (
            "Deeply analyze a PowerPoint (.pptx) deck: extract text, optionally "
            "render slides for visual analysis, and produce per-slide structured "
            "insights (risks, decisions, actions) plus a deck-level executive "
            "summary. Saves full analysis as JSON for follow-up Q&A."
        )

    @property
    def parameters(self) -> dict[str, Any]:
        return {
            "type": "object",
            "properties": {
                "path": {
                    "type": "string",
                    "description": "Path to the .pptx file to analyze.",
                },
                "model": {
                    "type": "string",
                    "description": (
                        "Vision-capable LLM model for analysis "
                        "(default from config or gpt-4o-mini)."
                    ),
                },
                "output_path": {
                    "type": "string",
                    "description": (
                        "Path to save the analysis JSON "
                        "(default: {name}_analysis.json alongside the .pptx)."
                    ),
                },
            },
            "required": ["path"],
        }

    def _cache_slide_analysis(
        self, file_path: str, slide_number: int, analysis: dict[str, Any]
    ) -> str:
        """Cache a single slide's analysis and return the cache key."""
        if not self._cache:
            return ""
        full_json = json.dumps(analysis, ensure_ascii=False, default=str)
        return self._cache.store(
            "analyze_pptx",
            {"path": file_path, "slide": slide_number},
            full_json,
            "",
            token_estimate=len(full_json) // 4,
        )

    async def execute(  # type: ignore[override]
        self,
        path: str,
        model: str | None = None,
        output_path: str | None = None,
        **kwargs: Any,
    ) -> ToolResult:
        resolved_model = model or self._vision_model

        try:
            file_path = _resolve_path(str(path), self._workspace, self._allowed_dir)
        except PermissionError as exc:
            return ToolResult.fail(str(exc), error_type="permission_denied")

        if not file_path.exists():
            return ToolResult.fail(f"File not found: {path}", error_type="not_found")
        if not file_path.is_file():
            return ToolResult.fail(f"Not a file: {path}", error_type="invalid_path")
        if file_path.suffix.lower() != ".pptx":
            return ToolResult.fail(f"Not a .pptx file: {path}", error_type="invalid_format")

        # Level 1: Extract text
        try:
            slides_data = _extract_slides_data(file_path)
        except ImportError as e:
            return ToolResult.fail(str(e), error_type="missing_dependency")
        except Exception as e:
            return ToolResult.fail(
                f"Error reading PowerPoint file: {e}", error_type="extraction_error"
            )

        if not slides_data:
            return ToolResult.ok("The presentation has no slides.")

        # Level 2: Render slides (optional — needs soffice + pdftoppm)
        render_dir = Path(tempfile.mkdtemp(prefix="pptx_render_"))
        try:
            images = await _render_slides(file_path, render_dir)
        except Exception:  # crash-barrier: rendering is optional
            images = None
        mode = "vision" if images else "text-only"

        # Map rendered images to slide numbers
        image_map: dict[int, Path] = {}
        if images:
            for i, img in enumerate(images):
                slide_num = i + 1
                if slide_num <= len(slides_data):
                    image_map[slide_num] = img

        # Level 2+3: Per-slide LLM analysis (parallel with concurrency limit)
        sem = asyncio.Semaphore(MAX_CONCURRENT)

        async def _bounded(sd: dict[str, Any]) -> dict[str, Any]:
            async with sem:
                img = image_map.get(sd["slide_number"])
                return await _analyze_slide(sd, img, resolved_model)

        try:
            slide_analyses = list(await asyncio.gather(*[_bounded(sd) for sd in slides_data]))
        except Exception as e:
            return ToolResult.fail(f"Error during slide analysis: {e}", error_type="analysis_error")

        # Cache each slide analysis
        for sa in slide_analyses:
            self._cache_slide_analysis(str(file_path), sa.get("slide_number", 0), sa)

        # Level 3: Deck synthesis
        try:
            synthesis = await _synthesize_deck(slide_analyses, resolved_model)
        except Exception as e:
            return ToolResult.fail(f"Error during deck synthesis: {e}", error_type="analysis_error")

        # Save full analysis JSON
        if output_path:
            try:
                out = _resolve_path(output_path, self._workspace, self._allowed_dir)
            except PermissionError as exc:
                return ToolResult.fail(str(exc), error_type="permission_denied")
        else:
            out = file_path.parent / f"{file_path.stem}_analysis.json"

        full_analysis = {
            "file": file_path.name,
            "total_slides": len(slides_data),
            "analysis_mode": mode,
            "model": resolved_model,
            **synthesis,
            "slides": slide_analyses,
        }
        out.parent.mkdir(parents=True, exist_ok=True)
        out.write_text(
            json.dumps(full_analysis, ensure_ascii=False, indent=2),
            encoding="utf-8",
        )

        # Format readable summary
        summary = _format_analysis_output(file_path.name, len(slides_data), mode, synthesis, out)
        return ToolResult.ok(summary)


# ---------------------------------------------------------------------------
# Tool: pptx_get_slide (Level 4 helper)
# ---------------------------------------------------------------------------


class PptxGetSlideTool(Tool):
    """Retrieve cached slide data from a previous read_pptx or analyze_pptx call."""

    readonly = True
    cacheable = False

    def __init__(self, cache: ToolResultCache) -> None:
        self._cache = cache

    @property
    def name(self) -> str:
        return "pptx_get_slide"

    @property
    def description(self) -> str:
        return (
            "Retrieve the full content of a specific slide from a previously "
            "cached read_pptx or analyze_pptx result. Use the cache_key from "
            "the read_pptx or analyze_pptx output."
        )

    @property
    def parameters(self) -> dict[str, Any]:
        return {
            "type": "object",
            "properties": {
                "cache_key": {
                    "type": "string",
                    "description": "The cache key from a prior read_pptx or analyze_pptx call.",
                },
            },
            "required": ["cache_key"],
        }

    async def execute(  # type: ignore[override]
        self,
        cache_key: str,
        **kwargs: Any,
    ) -> ToolResult:
        entry = self._cache.get(cache_key)
        if entry is None:
            return ToolResult.fail(
                f"No cached data found for key: {cache_key}",
                error_type="not_found",
            )
        return ToolResult.ok(entry.full_output)
