"""Tests for PowerPoint tools (ReadPptxTool, AnalyzePptxTool, PptxGetSlideTool)."""

from __future__ import annotations

import json
from pathlib import Path
from unittest.mock import AsyncMock, MagicMock, patch

import pytest

from nanobot.agent.tools.powerpoint import (
    AnalyzePptxTool,
    PptxGetSlideTool,
    ReadPptxTool,
    _extract_slides_data,
    _format_slides_markdown,
    _parse_json_response,
)
from nanobot.agent.tools.result_cache import ToolResultCache

pptx = pytest.importorskip("pptx")


# ---------------------------------------------------------------------------
# Fixtures
# ---------------------------------------------------------------------------


@pytest.fixture()
def sample_pptx(tmp_path: Path) -> Path:
    """Create a small .pptx file for testing."""
    from pptx import Presentation

    prs = Presentation()
    # Slide 1: title + body
    layout = prs.slide_layouts[1]  # title + content layout
    slide1 = prs.slides.add_slide(layout)
    slide1.shapes.title.text = "Project Overview"
    slide1.placeholders[1].text = "This is the project overview slide.\nKey milestone: Q3 launch."

    # Slide 2: just a title
    layout2 = prs.slide_layouts[0]  # title slide
    slide2 = prs.slides.add_slide(layout2)
    slide2.shapes.title.text = "Budget Summary"
    slide2.placeholders[1].text = "Total budget: $2M"

    # Slide 3: with speaker notes
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Timeline"
    slide3.placeholders[1].text = "Phase 1: Jan-Mar\nPhase 2: Apr-Jun"
    notes_slide = slide3.notes_slide
    notes_slide.notes_text_frame.text = "Remember to discuss dependencies."

    path = tmp_path / "deck.pptx"
    prs.save(str(path))
    return path


@pytest.fixture()
def sample_pptx_with_table(tmp_path: Path) -> Path:
    """Create a .pptx file with a table."""
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank layout
    table = slide.shapes.add_table(3, 2, Inches(1), Inches(1), Inches(4), Inches(2)).table
    table.cell(0, 0).text = "Name"
    table.cell(0, 1).text = "Role"
    table.cell(1, 0).text = "Alice"
    table.cell(1, 1).text = "PM"
    table.cell(2, 0).text = "Bob"
    table.cell(2, 1).text = "Dev"

    path = tmp_path / "table_deck.pptx"
    prs.save(str(path))
    return path


@pytest.fixture()
def cache(tmp_path: Path) -> ToolResultCache:
    return ToolResultCache(workspace=tmp_path)


@pytest.fixture()
def read_tool(tmp_path: Path, cache: ToolResultCache) -> ReadPptxTool:
    return ReadPptxTool(workspace=tmp_path, cache=cache)


@pytest.fixture()
def analyze_tool(tmp_path: Path, cache: ToolResultCache) -> AnalyzePptxTool:
    return AnalyzePptxTool(workspace=tmp_path, cache=cache, vision_model="test-model")


@pytest.fixture()
def get_slide_tool(cache: ToolResultCache) -> PptxGetSlideTool:
    return PptxGetSlideTool(cache=cache)


# ---------------------------------------------------------------------------
# Helper tests
# ---------------------------------------------------------------------------


def test_parse_json_response_clean() -> None:
    raw = '{"title": "Test", "summary": "A test"}'
    result = _parse_json_response(raw)
    assert result["title"] == "Test"


def test_parse_json_response_with_code_fence() -> None:
    raw = '```json\n{"title": "Test"}\n```'
    result = _parse_json_response(raw)
    assert result["title"] == "Test"


def test_parse_json_response_invalid() -> None:
    result = _parse_json_response("not json at all")
    assert "raw_response" in result


def test_extract_slides_data(sample_pptx: Path) -> None:
    slides = _extract_slides_data(sample_pptx)
    assert len(slides) == 3
    assert slides[0]["slide_number"] == 1
    assert "Project Overview" in slides[0]["title"]
    assert any("milestone" in b.lower() for b in slides[0]["text_blocks"])


def test_extract_slides_with_notes(sample_pptx: Path) -> None:
    slides = _extract_slides_data(sample_pptx)
    slide3 = slides[2]
    assert "dependencies" in slide3["notes"].lower()


def test_extract_slides_with_table(sample_pptx_with_table: Path) -> None:
    slides = _extract_slides_data(sample_pptx_with_table)
    assert len(slides) == 1
    tables = slides[0]["tables"]
    assert len(tables) == 1
    assert tables[0][0] == ["Name", "Role"]
    assert tables[0][1] == ["Alice", "PM"]


def test_format_slides_markdown(sample_pptx: Path) -> None:
    slides = _extract_slides_data(sample_pptx)
    md = _format_slides_markdown("deck.pptx", slides)
    assert "# PowerPoint: deck.pptx (3 slides)" in md
    assert "## Slide 1" in md
    assert "## Slide 2" in md
    assert "## Slide 3" in md
    assert "Project Overview" in md


# ---------------------------------------------------------------------------
# ReadPptxTool tests
# ---------------------------------------------------------------------------


async def test_read_pptx_happy_path(
    read_tool: ReadPptxTool,
    sample_pptx: Path,
    cache: ToolResultCache,
) -> None:
    result = await read_tool.execute(path=str(sample_pptx))
    assert result.success
    data = json.loads(result.output)
    assert data["total_slides"] == 3
    assert data["file"] == "deck.pptx"
    # Each slide should have a cache_key
    for slide in data["slides"]:
        assert "cache_key" in slide
        assert "slide_number" in slide
    # Verify slide 1 metadata
    assert "Project Overview" in data["slides"][0]["title"]


async def test_read_pptx_cache_contains_slide_data(
    read_tool: ReadPptxTool,
    sample_pptx: Path,
    cache: ToolResultCache,
) -> None:
    result = await read_tool.execute(path=str(sample_pptx))
    data = json.loads(result.output)
    cache_key = data["slides"][0]["cache_key"]
    entry = cache.get(cache_key)
    assert entry is not None
    cached = json.loads(entry.full_output)
    assert cached["slide_number"] == 1
    assert "text_blocks" in cached


async def test_read_pptx_with_tables(
    tmp_path: Path,
    sample_pptx_with_table: Path,
    cache: ToolResultCache,
) -> None:
    tool = ReadPptxTool(workspace=tmp_path, cache=cache)
    result = await tool.execute(path=str(sample_pptx_with_table))
    assert result.success
    data = json.loads(result.output)
    assert data["slides"][0]["has_tables"] is True


async def test_read_pptx_with_notes(
    read_tool: ReadPptxTool,
    sample_pptx: Path,
) -> None:
    result = await read_tool.execute(path=str(sample_pptx))
    data = json.loads(result.output)
    # Slide 3 has notes
    assert data["slides"][2]["has_notes"] is True
    # Slides 1 and 2 do not
    assert data["slides"][0]["has_notes"] is False


async def test_read_pptx_file_not_found(read_tool: ReadPptxTool, tmp_path: Path) -> None:
    result = await read_tool.execute(path=str(tmp_path / "missing.pptx"))
    assert not result.success
    assert "not found" in result.output.lower()


async def test_read_pptx_wrong_extension(read_tool: ReadPptxTool, tmp_path: Path) -> None:
    txt = tmp_path / "data.txt"
    txt.write_text("not a pptx")
    result = await read_tool.execute(path=str(txt))
    assert not result.success
    assert ".pptx" in result.output.lower()


async def test_read_pptx_not_a_file(read_tool: ReadPptxTool, tmp_path: Path) -> None:
    d = tmp_path / "some_dir.pptx"
    d.mkdir()
    result = await read_tool.execute(path=str(d))
    assert not result.success
    assert "not a file" in result.output.lower()


async def test_read_pptx_path_traversal_blocked(
    tmp_path: Path,
    sample_pptx: Path,
) -> None:
    """Tool with allowed_dir should block paths outside the workspace."""
    safe = tmp_path / "safe"
    safe.mkdir()
    tool = ReadPptxTool(workspace=tmp_path, allowed_dir=safe)
    result = await tool.execute(path=str(sample_pptx))
    assert not result.success
    assert "outside" in result.output.lower()


async def test_read_pptx_no_cache(tmp_path: Path, sample_pptx: Path) -> None:
    """Tool works without cache — no cache_key in metadata."""
    tool = ReadPptxTool(workspace=tmp_path, cache=None)
    result = await tool.execute(path=str(sample_pptx))
    assert result.success
    data = json.loads(result.output)
    assert data["total_slides"] == 3
    # No cache_key when cache is None
    assert "cache_key" not in data["slides"][0]


# ---------------------------------------------------------------------------
# AnalyzePptxTool tests
# ---------------------------------------------------------------------------


def _mock_llm_response(content: dict) -> AsyncMock:
    """Create a mock litellm.acompletion that returns structured JSON."""
    mock_response = MagicMock()
    mock_response.choices = [MagicMock()]
    mock_response.choices[0].message.content = json.dumps(content)
    return AsyncMock(return_value=mock_response)


_SLIDE_ANALYSIS = {
    "title": "Test Slide",
    "summary": "A test slide about testing.",
    "key_points": ["Testing is important"],
    "risks": ["Risk of not testing"],
}

_DECK_SYNTHESIS = {
    "executive_summary": "This deck is about testing.",
    "risks": ["Risk of not testing (slide 1)"],
    "decisions": [],
    "action_items": ["Write more tests"],
    "deadlines": [],
    "unanswered_questions": [],
    "themes": ["Testing"],
}


async def test_analyze_pptx_text_only(
    analyze_tool: AnalyzePptxTool,
    sample_pptx: Path,
    tmp_path: Path,
) -> None:
    """Analyze in text-only mode (no LibreOffice available)."""
    with (
        patch("nanobot.agent.tools.powerpoint._call_llm") as mock_llm,
        patch("nanobot.agent.tools.powerpoint._render_slides", return_value=None),
    ):
        # Mock the per-slide and synthesis calls
        slide_result = _SLIDE_ANALYSIS.copy()
        synthesis_result = _DECK_SYNTHESIS.copy()

        call_idx = 0

        async def _mock_call(messages, model):
            nonlocal call_idx
            call_idx += 1
            if call_idx <= 3:
                return json.dumps(slide_result)
            return json.dumps(synthesis_result)

        mock_llm.side_effect = _mock_call

        result = await analyze_tool.execute(path=str(sample_pptx))
        assert result.success
        assert "text-only" in result.output.lower()
        assert "executive summary" in result.output.lower()
        assert "Testing" in result.output or "testing" in result.output

        # Verify analysis JSON was saved
        output_json = sample_pptx.parent / "deck_analysis.json"
        assert output_json.exists()
        analysis = json.loads(output_json.read_text(encoding="utf-8"))
        assert analysis["total_slides"] == 3
        assert analysis["analysis_mode"] == "text-only"


async def test_analyze_pptx_custom_output_path(
    analyze_tool: AnalyzePptxTool,
    sample_pptx: Path,
    tmp_path: Path,
) -> None:
    """Verify custom output_path for the analysis JSON."""
    custom_out = tmp_path / "custom" / "result.json"

    async def _mock_call(messages, model):
        return json.dumps(_SLIDE_ANALYSIS)

    with (
        patch("nanobot.agent.tools.powerpoint._call_llm", side_effect=_mock_call),
        patch("nanobot.agent.tools.powerpoint._render_slides", return_value=None),
        patch(
            "nanobot.agent.tools.powerpoint._synthesize_deck",
            return_value=_DECK_SYNTHESIS,
        ),
    ):
        result = await analyze_tool.execute(path=str(sample_pptx), output_path=str(custom_out))
        assert result.success
        assert custom_out.exists()


async def test_analyze_pptx_file_not_found(
    analyze_tool: AnalyzePptxTool,
    tmp_path: Path,
) -> None:
    result = await analyze_tool.execute(path=str(tmp_path / "missing.pptx"))
    assert not result.success
    assert "not found" in result.output.lower()


async def test_analyze_pptx_wrong_extension(
    analyze_tool: AnalyzePptxTool,
    tmp_path: Path,
) -> None:
    txt = tmp_path / "data.txt"
    txt.write_text("not pptx")
    result = await analyze_tool.execute(path=str(txt))
    assert not result.success
    assert ".pptx" in result.output.lower()


# ---------------------------------------------------------------------------
# PptxGetSlideTool tests
# ---------------------------------------------------------------------------


async def test_get_slide_from_cache(
    read_tool: ReadPptxTool,
    get_slide_tool: PptxGetSlideTool,
    sample_pptx: Path,
) -> None:
    """Read a pptx, then retrieve a cached slide."""
    read_result = await read_tool.execute(path=str(sample_pptx))
    data = json.loads(read_result.output)
    cache_key = data["slides"][0]["cache_key"]

    result = await get_slide_tool.execute(cache_key=cache_key)
    assert result.success
    slide_data = json.loads(result.output)
    assert slide_data["slide_number"] == 1
    assert "text_blocks" in slide_data


async def test_get_slide_missing_key(get_slide_tool: PptxGetSlideTool) -> None:
    result = await get_slide_tool.execute(cache_key="nonexistent_key_12345")
    assert not result.success
    assert "no cached data" in result.output.lower()


async def test_get_slide_all_slides_cached(
    read_tool: ReadPptxTool,
    get_slide_tool: PptxGetSlideTool,
    sample_pptx: Path,
) -> None:
    """All slides from a read_pptx call should be individually retrievable."""
    read_result = await read_tool.execute(path=str(sample_pptx))
    data = json.loads(read_result.output)

    for slide_meta in data["slides"]:
        result = await get_slide_tool.execute(cache_key=slide_meta["cache_key"])
        assert result.success
        slide_data = json.loads(result.output)
        assert slide_data["slide_number"] == slide_meta["slide_number"]
